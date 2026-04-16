"""
generate_slides.py
Generates final/CampusShare_Storage_Slides.pptx from slide content.
Run: python3 scripts/generate_slides.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os

# ── Color palette ────────────────────────────────────────────────────────────
NAVY   = RGBColor(0x1B, 0x3A, 0x5C)
ORANGE = RGBColor(0xE8, 0x77, 0x22)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
GRAY_BG = RGBColor(0xF2, 0xF2, 0xF2)
GRAY_BORDER = RGBColor(0xAA, 0xAA, 0xAA)
BLACK  = RGBColor(0x1A, 0x1A, 0x1A)
LIGHT_NAVY = RGBColor(0x2E, 0x5F, 0x9E)

# ── Slide size: 16:9 ─────────────────────────────────────────────────────────
SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

# ── Layout constants ──────────────────────────────────────────────────────────
MARGIN = Inches(0.6)
CONTENT_TOP = Inches(1.55)
CONTENT_H   = Inches(5.5)
CONTENT_W   = SLIDE_W - 2 * MARGIN


def new_prs():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs):
    blank_layout = prs.slide_layouts[6]  # truly blank
    return prs.slides.add_slide(blank_layout)


# ── Drawing helpers ───────────────────────────────────────────────────────────

def add_rect(slide, left, top, width, height, fill_color=None, line_color=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    fill = shape.fill
    if fill_color:
        fill.solid()
        fill.fore_color.rgb = fill_color
    else:
        fill.background()
    line = shape.line
    if line_color:
        line.color.rgb = line_color
        line.width = Pt(1.2)
    else:
        line.fill.background()
    return shape


def set_para(para, text, font_size, bold=False, color=BLACK, align=PP_ALIGN.LEFT, space_before=0):
    para.clear()
    run = para.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Calibri"
    para.alignment = align
    para.space_before = Pt(space_before)


def add_textbox(slide, left, top, width, height, text, font_size,
                bold=False, color=BLACK, align=PP_ALIGN.LEFT,
                word_wrap=True, v_anchor=None):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    if v_anchor:
        from pptx.enum.text import MSO_ANCHOR
        tf.auto_size = None
    para = tf.paragraphs[0]
    set_para(para, text, font_size, bold=bold, color=color, align=align)
    return txBox


def add_header_bar(slide, title_text, subtitle_text=None):
    """Navy bar across the top with white title."""
    bar = add_rect(slide, 0, 0, SLIDE_W, Inches(1.3), fill_color=NAVY)
    # Title text
    txBox = slide.shapes.add_textbox(MARGIN, Inches(0.12), SLIDE_W - 2*MARGIN, Inches(1.05))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    set_para(p, title_text, 28, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle_text:
        p2 = tf.add_paragraph()
        set_para(p2, subtitle_text, 14, bold=False, color=RGBColor(0xCC, 0xD9, 0xEA), align=PP_ALIGN.LEFT)
    return bar


def add_bullet_points(slide, bullets, left=None, top=None, width=None, height=None,
                       font_size=18, color=BLACK, indent_first=False):
    """Add a list of (text, is_bold, is_sub) tuples as bullet points."""
    if left is None:  left = MARGIN
    if top is None:   top = CONTENT_TOP
    if width is None: width = CONTENT_W
    if height is None: height = CONTENT_H

    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    first = True
    for item in bullets:
        if isinstance(item, str):
            text, bold, sub = item, False, False
        elif len(item) == 2:
            text, bold = item; sub = False
        else:
            text, bold, sub = item

        if first:
            para = tf.paragraphs[0]
            first = False
        else:
            para = tf.add_paragraph()

        run = para.add_run()
        run.text = ("  " if sub else "• ") + text if not sub else "    ◦ " + text
        run.font.size = Pt(font_size - 2 if sub else font_size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = "Calibri"
        para.space_before = Pt(4 if sub else 8)
        para.alignment = PP_ALIGN.LEFT

    return txBox


def add_speaker_notes(slide, notes_text):
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.text = notes_text


def add_placeholder_box(slide, left, top, width, height, label):
    """Gray dashed-border box for a diagram that isn't ready yet."""
    box = add_rect(slide, left, top, width, height,
                   fill_color=GRAY_BG, line_color=GRAY_BORDER)
    # label
    txBox = slide.shapes.add_textbox(
        left + Inches(0.1), top + height/2 - Inches(0.3),
        width - Inches(0.2), Inches(0.6)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    set_para(p, label, 14, bold=False,
             color=GRAY_BORDER, align=PP_ALIGN.CENTER)
    return box


def add_orange_accent(slide, text, left, top, width=None, height=Inches(0.5)):
    """Highlighted stat callout in orange."""
    if width is None: width = Inches(5)
    box = add_rect(slide, left, top, width, height, fill_color=ORANGE)
    txBox = slide.shapes.add_textbox(
        left + Inches(0.1), top, width - Inches(0.2), height
    )
    tf = txBox.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    set_para(p, text, 16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    return txBox


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDES
# ═══════════════════════════════════════════════════════════════════════════════

def slide_01_title(prs):
    slide = blank_slide(prs)
    # Full navy background
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill_color=NAVY)

    # Decorative orange accent bar
    add_rect(slide, 0, Inches(3.9), SLIDE_W, Inches(0.08), fill_color=ORANGE)

    # Main title
    txBox = slide.shapes.add_textbox(MARGIN, Inches(1.2), SLIDE_W - 2*MARGIN, Inches(1.8))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    set_para(p, "CampusShare Storage", 44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Subtitle
    txBox2 = slide.shapes.add_textbox(MARGIN, Inches(2.9), SLIDE_W - 2*MARGIN, Inches(0.9))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    set_para(p2, "The Campus P2P Storage Marketplace", 24, bold=False,
             color=RGBColor(0xCC, 0xD9, 0xEA), align=PP_ALIGN.CENTER)

    # Course + group info
    txBox3 = slide.shapes.add_textbox(MARGIN, Inches(4.2), SLIDE_W - 2*MARGIN, Inches(1.8))
    tf3 = txBox3.text_frame
    tf3.word_wrap = True
    for line in [
        ("MIDS301 Introduction to Information — Spring 2026", 16, False),
        ("Group members: [insert names]", 15, False),
        ("Presentation date: [insert date]", 15, False),
    ]:
        if tf3.paragraphs[0].runs:
            p = tf3.add_paragraph()
        else:
            p = tf3.paragraphs[0]
        set_para(p, line[0], line[1], bold=line[2],
                 color=RGBColor(0xCC, 0xD9, 0xEA), align=PP_ALIGN.CENTER)
        p.space_before = Pt(6)

    add_speaker_notes(slide,
        "Good morning everyone. We're presenting CampusShare Storage — a peer-to-peer marketplace "
        "that lets university students store belongings with each other, right on campus. Over the "
        "next ten minutes we'll walk you through our strategy, data model, and process design.")
    return slide


def slide_02_problem(prs):
    slide = blank_slide(prs)
    add_header_bar(slide, "The Problem",
                   "1.1 million international students. Nowhere to store their things.")

    bullets = [
        ("1,126,690 international students enrolled in U.S. universities (IIE, 2024)", True, False),
        ("Commercial storage averages $150/month — miles from campus", False, False),
        ("No car, no truck, no local family — access is the real barrier", False, False),
        ("Summer housing ends; belongings have no place to go", False, False),
    ]
    add_bullet_points(slide, bullets, font_size=20)

    # Stat callout
    add_orange_accent(slide, "$150/month  |  Miles away  |  No vehicle",
                      left=Inches(3.0), top=Inches(6.1), width=Inches(7.3))

    add_speaker_notes(slide,
        "Every year, over a million international students face the same crisis — their lease ends "
        "in May but their flight home isn't until June, and they have no practical place to store "
        "three months' worth of belongings. Commercial storage exists, but it's expensive, it "
        "requires a vehicle, and it's designed for suburban homeowners, not carless college students.")
    return slide


def slide_03_solution(prs):
    slide = blank_slide(prs)
    add_header_bar(slide, "Our Solution",
                   ".edu-verified peer storage, walking distance, no truck needed")

    bullets = [
        ("Students with unused space become Hosts and earn passive income", False, False),
        ("Students needing storage become Renters — book online in minutes", False, False),
        ("Both sides verified with .edu email — same campus community", True, False),
        ("Flexible terms: 1 week to 3 months", False, False),
        ("Physical handover happens on campus, on foot", False, False),
    ]
    add_bullet_points(slide, bullets, font_size=20)

    # Simple flow diagram (text-based)
    add_orange_accent(slide, "Renter  ←→  Platform (.edu verified)  ←→  Host",
                      left=Inches(2.5), top=Inches(6.1), width=Inches(8.3))

    add_speaker_notes(slide,
        "CampusShare connects these two groups that already live together. A Renter submits a "
        "booking request; a Host accepts. Both are verified members of the same campus. The item "
        "handover happens at a mutually agreed on-campus location — no truck, no drive, no storage "
        "facility account required.")
    return slide


def slide_04_market(prs):
    slide = blank_slide(prs)
    add_header_bar(slide, "Market & Customer",
                   "$1.5M addressable GMV per campus — at one mid-sized university")

    bullets = [
        ("Primary: international students — summer storage need is structural", True, False),
        ("Secondary: domestic students between leases or on internship", False, False),
        ("Scope: one university at launch (single-campus, by design)", False, False),
        ("5,000 international students × avg. $300 stored goods = $1.5M GMV", True, False),
        ("CampusShare at 10% commission = $150K annual revenue per campus", True, False),
    ]
    add_bullet_points(slide, bullets, font_size=19)

    # Unit economics box
    box = add_rect(slide, Inches(3.0), Inches(5.9), Inches(7.3), Inches(0.85),
                   fill_color=RGBColor(0xE8, 0xF0, 0xF8), line_color=LIGHT_NAVY)
    txBox = slide.shapes.add_textbox(Inches(3.1), Inches(5.95), Inches(7.1), Inches(0.75))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    set_para(p, "5,000 students × $300 = $1.5M GMV  →  $150K revenue @ 10%",
             15, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

    add_speaker_notes(slide,
        "The market is real and calculable. At a mid-sized university with 5,000 international "
        "students, if each stores an average of $300 worth of goods over the summer, the addressable "
        "gross merchandise value exceeds $1.5 million per campus per year. CampusShare takes 10%, "
        "which is $150,000 in annual revenue from a single campus — before any insurance add-on.")
    return slide


def slide_05_porter(prs):
    slide = blank_slide(prs)
    add_header_bar(slide, "Porter's Five Forces",
                   "Highest risk: substitutes.  Biggest moat: .edu trust + network effects")

    # Left: bullets (narrower to leave room for diagram placeholder)
    bullets = [
        ("Threat of New Entrants — Moderate: low tech barrier, but .edu trust wall + review moat take time", False, False),
        ("Suppliers (Hosts) — Low–Moderate: review reputation creates soft lock-in", False, False),
        ("Buyers (Renters) — Moderate: proximity advantage dampens switching", False, False),
        ("Substitutes — HIGH ⚠: U-Haul, Public Storage, friend storage, shipping home", True, False),
        ("Rivalry — Low–Moderate: no .edu-gated campus P2P storage competitor today", False, False),
    ]
    add_bullet_points(slide, bullets,
                      left=MARGIN, top=CONTENT_TOP,
                      width=Inches(7.0), height=CONTENT_H, font_size=17)

    # Right: diagram placeholder
    add_placeholder_box(slide,
                        left=Inches(7.8), top=CONTENT_TOP,
                        width=Inches(5.0), height=Inches(5.2),
                        label="[ Insert porter_five_forces.png here ]")

    add_speaker_notes(slide,
        "The most serious force is substitution — students have real alternatives like U-Haul or "
        "asking a friend. But those options require a vehicle, trust in a stranger, or large costs. "
        "No current competitor sits at the intersection of verified student identity, campus proximity, "
        "and short-term peer storage. That's our whitespace.")
    return slide


def slide_06_strategy(prs):
    slide = blank_slide(prs)
    add_header_bar(slide, "Generic Strategy",
                   "Focused Differentiation — not cheap, not broad, deeply trusted")

    bullets = [
        ("Focus: one campus, one verified community — not the whole storage market", True, False),
        ("Differentiation: .edu identity + walking distance + photo documentation + dual reviews", True, False),
        ("Not Cost Leadership: can't out-scale U-Haul on unit cost as a new entrant", False, False),
        ("Not Broad Differentiation: expanding scope before PMF = stuck in the middle", False, False),
        ('Porter (1985): attempting both without focus risks losing to everyone', False, False),
    ]
    add_bullet_points(slide, bullets, font_size=19)

    add_orange_accent(slide, "Position: Focused Differentiation (narrow scope + high differentiation)",
                      left=Inches(2.0), top=Inches(6.1), width=Inches(9.3))

    add_speaker_notes(slide,
        "Porter warns against being 'stuck in the middle' — trying to be cheap and differentiated "
        "without focus. We've deliberately chosen the bottom-right quadrant: narrow scope, high "
        "differentiation. The .edu wall, the review system, the on-campus proximity — these create "
        "a value proposition that a general marketplace can't replicate and a big storage company "
        "has no reason to build.")
    return slide


def slide_07_revenue(prs):
    slide = blank_slide(prs)
    add_header_bar(slide, "Revenue Model",
                   "10% commission + optional insurance — simple, aligned, scalable")

    bullets = [
        ("Primary: 10% platform commission on each completed booking (auto-deducted)", True, False),
        ("Secondary: optional insurance add-on, $5–$15 per booking (tiered by declared value)", True, False),
        ("Airbnb comparison: host ~3% + guest 14–16% = far higher take rate (Airbnb, 2023)", False, False),
        ("No subscriptions — conflicts with 'save money' positioning", False, False),
        ("No ads — user base too small at launch to monetize attention", False, False),
    ]
    add_bullet_points(slide, bullets, font_size=19)

    add_speaker_notes(slide,
        "Revenue is simple and per-transaction. We take 10% when a booking completes — no payment "
        "until value is delivered. The insurance add-on is optional at checkout; it covers the "
        "Renter's declared item value for loss or damage. We deliberately excluded subscriptions "
        "and advertising because they'd either conflict with the student value prop or require "
        "scale we don't yet have.")
    return slide


def slide_08_erd(prs):
    slide = blank_slide(prs)
    add_header_bar(slide, "Information Design — ERD",
                   "9 entities, 2 associative entities — Booking is the hub")

    # Left: bullets
    bullets = [
        ("Student is one entity — same person can be Host or Renter", True, False),
        ("3 functional layers: actor identity / core transactions / trust & risk", False, False),
        ("Booking resolves Student (Renter) × StorageListing M:N — carries status lifecycle", True, False),
        ("Review resolves mutual peer evaluation M:N — tied to verified Bookings", True, False),
        ("Downstream: StoredItem, Payment, Dispute, InsurancePolicy all attach to Booking", False, False),
    ]
    add_bullet_points(slide, bullets,
                      left=MARGIN, top=CONTENT_TOP,
                      width=Inches(6.6), height=CONTENT_H, font_size=17)

    # Right: ERD placeholder
    add_placeholder_box(slide,
                        left=Inches(7.4), top=CONTENT_TOP,
                        width=Inches(5.4), height=Inches(5.2),
                        label="[ Insert ERD.png here ]")

    add_speaker_notes(slide,
        "The data model has nine entities organized in three layers. Everything flows through "
        "Booking — it's the central hub that connects the Renter, the Listing, and all downstream "
        "records. We modeled Host and Renter as roles of a single Student entity rather than "
        "separate entities, because the same person can flip between both roles over time. Booking "
        "and Review are associative entities — they resolve many-to-many relationships and carry "
        "their own attributes.")
    return slide


def slide_09_value_chain(prs):
    slide = blank_slide(prs)
    add_header_bar(slide, "Value Chain",
                   "Every primary activity is scoped to the campus boundary — by design")

    # Left: bullets
    bullets = [
        ("Inbound Logistics: Host onboarding, .edu verification, Renter registration", False, False),
        ("Operations: search ranking, booking state mgmt, escrow, item photo registry", False, False),
        ("Outbound Logistics: on-campus handover, timestamped confirmation, return", False, False),
        ("Marketing & Sales: student orgs, dorm boards, intl. student office, .edu lists", False, False),
        ("Service: dual reviews, dispute mediation, insurance claim assistance", False, False),
    ]
    add_bullet_points(slide, bullets,
                      left=MARGIN, top=CONTENT_TOP,
                      width=Inches(6.6), height=CONTENT_H, font_size=17)

    # Right: Value chain placeholder
    add_placeholder_box(slide,
                        left=Inches(7.4), top=CONTENT_TOP,
                        width=Inches(5.4), height=Inches(5.2),
                        label="[ Insert value_chain.png here ]")

    add_speaker_notes(slide,
        "Porter's value chain shows how each operational stage reinforces our strategy. Inbound "
        "and Marketing use the .edu wall to keep both sides verified. Operations uses data-driven "
        "matching to reduce friction. Outbound keeps logistics on-campus — walkable. And Service "
        "builds the trust asset that makes each subsequent transaction easier. Every activity is "
        "bounded by the campus, which is what makes the model work.")
    return slide


def slide_10_bpmn(prs):
    slide = blank_slide(prs)
    add_header_bar(slide, "Business Execution — BPMN",
                   "14 activities, 2 sub-processes, one complete storage transaction lifecycle")

    # Left: bullets
    bullets = [
        ("Stage 1 — Search & Select: Renter logs in, .edu verified, queries listings, submits Booking", False, False),
        ("Stage 2 — Confirm: Host approves/declines (XOR); Renter pays; System generates agreement", False, False),
        ("Stage 3 — Handover: both parties coordinate; System records 'active' status", False, False),
        ("Stage 4 — Closeout: return confirmed; mutual reviews written to Review table", False, False),
        ("Sub-processes: Payment Processing (Stripe) + Dispute Handling (exception path)", True, False),
    ]
    add_bullet_points(slide, bullets,
                      left=MARGIN, top=CONTENT_TOP,
                      width=Inches(6.6), height=CONTENT_H, font_size=17)

    # Right: BPMN placeholder
    add_placeholder_box(slide,
                        left=Inches(7.4), top=CONTENT_TOP,
                        width=Inches(5.4), height=Inches(5.2),
                        label="[ Insert BPMN.png here ]")

    add_speaker_notes(slide,
        "The BPMN shows the full booking lifecycle in BPMN 2.0 using a Database Interaction Layout. "
        "Every step that touches data has a dashed arrow connecting it to the correct database "
        "entity. The two sub-processes encapsulate complexity: Payment Processing hides the Stripe "
        "API interaction, and Dispute Handling isolates the exception path so it doesn't clutter "
        "the main success flow.")
    return slide


def slide_11_ethics(prs):
    slide = blank_slide(prs)
    add_header_bar(slide, "Ethics, Limitations & Future",
                   "Inference risk is real — and we've designed against it")

    col_w = Inches(3.9)
    col_top = CONTENT_TOP
    col_h = Inches(5.0)

    # Column headers
    for i, (title, col_x) in enumerate([
        ("Ethics", MARGIN),
        ("Limitations", MARGIN + col_w + Inches(0.15)),
        ("Future", MARGIN + 2*(col_w + Inches(0.15))),
    ]):
        add_rect(slide, col_x, col_top, col_w, Inches(0.42), fill_color=NAVY)
        txb = slide.shapes.add_textbox(col_x + Inches(0.1), col_top + Inches(0.04),
                                        col_w - Inches(0.2), Inches(0.38))
        p = txb.text_frame.paragraphs[0]
        set_para(p, title, 16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    ethics_bullets = [
        "Combined exposure: verified identity + Host address + item photos",
        "Disproportionate harm to international students",
        "Mitigations: encrypted photos (purged 90 days), address withheld until confirmed, campus drop-off option",
    ]
    lim_bullets = [
        "Single-campus scope only",
        "Sparse early reviews (bootstrapping problem)",
        "Insurance workflow undefined",
        "No admin entity in ERD",
        "Prohibited items not formally encoded",
    ]
    future_bullets = [
        "Cross-campus consortium",
        "ML-based matching",
        "Dynamic pricing",
        "Logistics courier add-on",
        "Licensed insurance partnership",
    ]

    bullet_top = col_top + Inches(0.52)
    bullet_h = Inches(4.5)
    for bullets, col_x in [
        (ethics_bullets, MARGIN),
        (lim_bullets, MARGIN + col_w + Inches(0.15)),
        (future_bullets, MARGIN + 2*(col_w + Inches(0.15))),
    ]:
        txb = slide.shapes.add_textbox(col_x, bullet_top, col_w, bullet_h)
        tf = txb.text_frame
        tf.word_wrap = True
        first = True
        for b in bullets:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            run = p.add_run()
            run.text = "• " + b
            run.font.size = Pt(14)
            run.font.name = "Calibri"
            run.font.color.rgb = BLACK
            p.space_before = Pt(6)

    add_speaker_notes(slide,
        "The most distinctive ethical risk isn't any single data point — it's the combination. "
        "Name plus address plus item photos creates a profile that could enable targeted theft, "
        "especially for international students who may struggle to engage local law enforcement. "
        "We've mitigated this by withholding addresses until a booking is confirmed and automatically "
        "purging item photos after 90 days. As for the future — once this model works on one campus, "
        "a cross-campus consortium is the natural next move.")
    return slide


def slide_12_closing(prs):
    slide = blank_slide(prs)
    # Navy background
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill_color=NAVY)
    add_rect(slide, 0, Inches(1.2), SLIDE_W, Inches(0.06), fill_color=ORANGE)
    add_rect(slide, 0, Inches(6.3), SLIDE_W, Inches(0.06), fill_color=ORANGE)

    # Headline
    txb = slide.shapes.add_textbox(MARGIN, Inches(0.2), SLIDE_W - 2*MARGIN, Inches(0.95))
    tf = txb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    set_para(p, "One campus. One community. One trusted marketplace.",
             30, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Three moat boxes
    moat_data = [
        ("#1  .edu identity wall", "Only verified students can participate"),
        ("#2  Dual-direction reviews", "Trust compounds with each transaction"),
        ("#3  Deposit-escrow mechanism", "Aligned incentives for both parties"),
    ]
    box_w = Inches(3.6)
    box_h = Inches(1.6)
    box_top = Inches(1.45)
    for i, (title, desc) in enumerate(moat_data):
        bx = MARGIN + i * (box_w + Inches(0.2))
        add_rect(slide, bx, box_top, box_w, box_h,
                 fill_color=LIGHT_NAVY, line_color=ORANGE)
        txb2 = slide.shapes.add_textbox(bx + Inches(0.1), box_top + Inches(0.1),
                                         box_w - Inches(0.2), box_h - Inches(0.2))
        tf2 = txb2.text_frame
        tf2.word_wrap = True
        p1 = tf2.paragraphs[0]
        set_para(p1, title, 16, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)
        p2 = tf2.add_paragraph()
        set_para(p2, desc, 13, bold=False, color=WHITE, align=PP_ALIGN.CENTER, space_before=6)

    # Bottom stats
    stats_bullets = [
        ("$1.5M GMV addressable per campus at zero marginal infrastructure cost", False, False),
        ("First-mover on campus = brand + reviews + relationships hard to displace", False, False),
    ]
    add_bullet_points(slide, stats_bullets,
                      top=Inches(3.3), font_size=18,
                      color=RGBColor(0xCC, 0xD9, 0xEA))

    # Thank you
    txb3 = slide.shapes.add_textbox(MARGIN, Inches(6.4), SLIDE_W - 2*MARGIN, Inches(0.8))
    tf3 = txb3.text_frame
    p3 = tf3.paragraphs[0]
    set_para(p3, "Thank you — we're happy to take questions.",
             18, bold=False, color=WHITE, align=PP_ALIGN.CENTER)

    add_speaker_notes(slide,
        "CampusShare is not trying to replace U-Haul. We're building something U-Haul cannot build: "
        "a trusted peer community within a single campus, gated by student identity. The moat gets "
        "stronger with every transaction — more reviews, more trust, more Hosts onboarded. We "
        "believe this model works at one campus, and then scales campus by campus. Thank you — "
        "we're happy to take questions.")
    return slide


# ═══════════════════════════════════════════════════════════════════════════════
# MAIN
# ═══════════════════════════════════════════════════════════════════════════════

def main():
    project_root = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
    output_path = os.path.join(project_root, "final", "CampusShare_Storage_Slides.pptx")
    os.makedirs(os.path.dirname(output_path), exist_ok=True)

    prs = new_prs()

    print("Building slides...")
    slide_01_title(prs)
    print("  1/12 Title")
    slide_02_problem(prs)
    print("  2/12 The Problem")
    slide_03_solution(prs)
    print("  3/12 Our Solution")
    slide_04_market(prs)
    print("  4/12 Market & Customer")
    slide_05_porter(prs)
    print("  5/12 Porter's Five Forces")
    slide_06_strategy(prs)
    print("  6/12 Generic Strategy")
    slide_07_revenue(prs)
    print("  7/12 Revenue Model")
    slide_08_erd(prs)
    print("  8/12 ERD")
    slide_09_value_chain(prs)
    print("  9/12 Value Chain")
    slide_10_bpmn(prs)
    print(" 10/12 BPMN")
    slide_11_ethics(prs)
    print(" 11/12 Ethics / Limitations / Future")
    slide_12_closing(prs)
    print(" 12/12 Closing Ask")

    prs.save(output_path)
    print(f"\nSaved → {output_path}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
